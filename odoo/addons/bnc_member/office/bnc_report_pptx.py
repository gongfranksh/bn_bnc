from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import os, sys


class bnc_report_ppt(object):
#    def __init__(self, object):
# self.url = object['para_my_url']['bn_2dfire_function_api']
# self.connection = object['para_connection']

    def create_ppt(self, data):
        # create presentation with 1 slide ------
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        zp=[]
        zp_value=[]
        pro_data=sorted(data,key=lambda result:result[0])
        for item in pro_data:
            zp.append(item[0])
            zp_value.append(item[1])
        chart_data = ChartData()
        chart_data.categories = zp
        chart_data.add_series('Series 1', zp_value)

        # add chart to slide --------------------
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        file_name='chart-011.pptx'
        file_path = os.path.join(sys.path[0], file_name)
        prs.save(file_path)

        with open(file_path, 'rb') as fp:
            data = fp.read().encode('base64')
        return data
